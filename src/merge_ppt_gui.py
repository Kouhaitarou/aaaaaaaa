import os
import sys
import copy
from tkinter import (
    Tk,
    Toplevel,
    Listbox,
    Button,
    Frame,
    END,
    filedialog,
    Text,
    Scrollbar,
    VERTICAL,
    RIGHT,
    Y,
    LEFT,
    BOTH,
)
from tkinter import messagebox
from pptx import Presentation


def load_presentation(path):
    return Presentation(path)


def copy_slide(source, target_prs):
    layout = target_prs.slide_layouts[6]
    new_slide = target_prs.slides.add_slide(layout)
    for shape in source.shapes:
        new_shape = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    if source.has_notes_slide:
        notes = source.notes_slide.notes_text_frame.text
        new_slide.notes_slide.notes_text_frame.text = notes
    return new_slide


class App:
    def __init__(self, master):
        self.master = master
        master.title('PowerPoint Merge Tool')
        self.files = []

        self.listbox = Listbox(master, selectmode='extended', width=60)
        self.listbox.pack(side=LEFT, fill=BOTH, expand=True)

        scrollbar = Scrollbar(master, orient=VERTICAL, command=self.listbox.yview)
        scrollbar.pack(side=RIGHT, fill=Y)
        self.listbox.config(yscrollcommand=scrollbar.set)

        btn_frame = Frame(master)
        Button(btn_frame, text='Add Files', command=self.add_files).pack(fill='x')
        Button(btn_frame, text='Move Up', command=lambda: self.move(-1)).pack(fill='x')
        Button(btn_frame, text='Move Down', command=lambda: self.move(1)).pack(fill='x')
        Button(btn_frame, text='Edit TOC', command=self.edit_toc).pack(fill='x')
        Button(btn_frame, text='Merge', command=self.merge).pack(fill='x')
        btn_frame.pack(side=LEFT)

        self.toc_window = None
        self.toc_text = None

    def add_files(self):
        paths = filedialog.askopenfilenames(filetypes=[('PowerPoint', '*.pptx')])
        for p in paths:
            if p not in self.files:
                self.files.append(p)
                self.listbox.insert(END, os.path.basename(p))

    def move(self, direction):
        selections = list(self.listbox.curselection())
        if not selections:
            return
        indices = selections if direction < 0 else reversed(selections)
        for idx in indices:
            new_idx = idx + direction
            if 0 <= new_idx < self.listbox.size():
                self.files[idx], self.files[new_idx] = self.files[new_idx], self.files[idx]
                text = self.listbox.get(idx)
                self.listbox.delete(idx)
                self.listbox.insert(new_idx, text)
                self.listbox.selection_set(new_idx)

    def edit_toc(self):
        if self.toc_window:
            return
        self.toc_window = Toplevel(self.master)
        self.toc_window.title('Edit Table of Contents')
        self.toc_text = Text(self.toc_window, width=60, height=20)
        content = []
        for idx, path in enumerate(self.files, 1):
            prs = Presentation(path)
            content.append(f'{idx}. {os.path.basename(path)} - {len(prs.slides)} slides')
        self.toc_text.insert(1.0, '\n'.join(content))
        self.toc_text.pack()
        Button(self.toc_window, text='Close', command=self.close_toc).pack()
        self.toc_window.protocol('WM_DELETE_WINDOW', self.close_toc)

    def close_toc(self):
        if self.toc_window:
            self.toc_window.destroy()
            self.toc_window = None

    def merge(self):
        if not self.files:
            messagebox.showinfo('Info', 'No files selected')
            return
        output_path = filedialog.asksaveasfilename(defaultextension='.pptx', filetypes=[('PowerPoint', '*.pptx')])
        if not output_path:
            return
        result = Presentation()
        if self.toc_text:
            toc_text = self.toc_text.get(1.0, END).strip()
        else:
            toc_text = '\n'.join(f'{i+1}. {os.path.basename(p)}' for i, p in enumerate(self.files))
        layout = result.slide_layouts[1]
        slide = result.slides.add_slide(layout)
        slide.shapes.title.text = '目次'
        body = slide.placeholders[1].text_frame
        body.text = toc_text

        for path in self.files:
            src_prs = Presentation(path)
            for slide in src_prs.slides:
                copy_slide(slide, result)
        result.save(output_path)
        messagebox.showinfo('Info', f'Merged to {output_path}')


def main():
    root = Tk()
    app = App(root)
    root.mainloop()


if __name__ == '__main__':
    main()
